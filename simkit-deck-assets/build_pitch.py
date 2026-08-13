#!/usr/bin/env python3
"""Simkit.ai pitch deck — Notion theme, doodle personas, icon cards. No blue."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ASSETS = Path(__file__).resolve().parent
CURSOR = Path(
    "/Users/bbhavnasi/.cursor/projects/Users-bbhavnasi-Desktop-GW-Apps-concur-agent/assets"
)
DOODLES = ASSETS / "doodles"
OUT_LOCAL = ASSETS / "Simkit.ai.pitch.pptx"
OUT_DESKTOP = Path("/Users/bbhavnasi/Desktop/Simkit.ai.pitch.pptx")
OUT_DOWNLOADS = Path("/Users/bbhavnasi/Downloads/Simkit.ai.notion.pptx")

# Notion palette — no blue
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE = RGBColor(0xF7, 0xF6, 0xF3)  # Notion sidebar / card
CALL_YELLOW = RGBColor(0xFB, 0xF3, 0xDB)
CALL_ORANGE = RGBColor(0xFA, 0xEB, 0xDD)
CALL_GREEN = RGBColor(0xED, 0xF3, 0xEC)
CALL_PINK = RGBColor(0xFC, 0xEB, 0xEB)
CALL_GRAY = RGBColor(0xF1, 0xF1, 0xEF)
INK = RGBColor(0x37, 0x35, 0x2F)  # Notion default text
MUTED = RGBColor(0x78, 0x77, 0x74)
FAINT = RGBColor(0x9B, 0x9A, 0x97)
RULE = RGBColor(0xE9, 0xE9, 0xE7)
ACCENT = RGBColor(0x37, 0x35, 0x2F)  # black accent (no blue)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)
FONT = "Helvetica Neue"


def asset(*names: str) -> Path:
    bases = [
        CURSOR,
        ASSETS,
        ASSETS / "icons",
        ASSETS / "doodles",
        DOODLES,
        Path("/Users/bbhavnasi/Desktop/GW-Apps/0.GWVA/GWVA-CLI/cli/internal/simruntime/ui/public/personas"),
    ]
    for name in names:
        for base in bases:
            p = base / name
            if p.exists():
                return p
    raise FileNotFoundError(names)


def doodle(n: str) -> Path:
    return asset(f"doodle-{n}.jpg")


def icon(name: str) -> Path:
    return asset(f"icon-{name}.png")


def set_run(run, *, size=18, bold=False, color=INK, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_text(slide, left, top, width, height, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_lines(slide, left, top, width, height, lines, *, size=14, color=MUTED):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, color=color)
    return box


def fill_slide(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def hairline(slide, left, top, width):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RULE
    sh.line.fill.background()


def card(slide, left, top, width, height, *, fill=SURFACE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = RULE
    sh.line.width = Pt(1)
    try:
        sh.adjustments[0] = 0.1
    except Exception:
        pass
    return sh


def pill(slide, left, top, width, height, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    try:
        sh.adjustments[0] = 0.2
    except Exception:
        pass
    return sh


def add_pic(slide, path, left, top, width, height=None):
    if height is None:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def footer(slide, label: str):
    hairline(slide, MARGIN, Inches(7.05), SLIDE_W - 2 * MARGIN)
    add_text(slide, MARGIN, Inches(7.12), Inches(3), Inches(0.28), "SIMKIT", size=10, bold=True, color=FAINT)
    add_text(
        slide,
        SLIDE_W - MARGIN - Inches(1.0),
        Inches(7.12),
        Inches(1.0),
        Inches(0.28),
        label,
        size=10,
        color=FAINT,
        align=PP_ALIGN.RIGHT,
    )


def kicker(slide, text: str):
    add_text(slide, MARGIN, Inches(0.4), Inches(10), Inches(0.28), text, size=11, bold=True, color=MUTED)


def headline(slide, text: str, *, top=Inches(0.75), size=34, width=Inches(11.5)):
    add_text(slide, MARGIN, top, width, Inches(1.3), text, size=size, bold=True, color=INK)


def sub(slide, text: str, *, top=Inches(2.15), width=Inches(10.5)):
    add_text(slide, MARGIN, top, width, Inches(0.7), text, size=15, color=MUTED)


def doodle_chip(slide, n: str, left, top, size=Inches(0.95)):
    """Circular-ish persona chip via rounded rect + doodle."""
    bg = pill(slide, left, top, size, size, CALL_GRAY)
    # inset the doodle slightly
    inset = Inches(0.06)
    add_pic(slide, doodle(n), left + inset, top + inset, size - 2 * inset, size - 2 * inset)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ══════════════════════════════════════════════════════════
    # 01 TITLE — doodle crowd + catchy line
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    fill_slide(s, WHITE)

    # doodle constellation on the right
    layout = [
        ("01", 8.6, 1.2, 1.35),
        ("40", 10.2, 1.0, 1.5),
        ("12", 11.6, 1.4, 1.2),
        ("05", 8.4, 2.9, 1.25),
        ("18", 9.9, 2.7, 1.55),
        ("30", 11.55, 3.0, 1.25),
        ("08", 9.0, 4.5, 1.2),
        ("44", 10.5, 4.4, 1.4),
        ("15", 12.0, 4.6, 1.15),
    ]
    for n, x, y, sz in layout:
        add_pic(s, doodle(n), Inches(x), Inches(y), Inches(sz), Inches(sz))

    add_text(s, MARGIN, Inches(1.6), Inches(7), Inches(0.35), "SIMKIT.AI", size=13, bold=True, color=MUTED)
    add_text(
        s,
        MARGIN,
        Inches(2.1),
        Inches(7.5),
        Inches(1.6),
        "Ship voice agents\nwith 100% confidence.",
        size=42,
        bold=True,
        color=INK,
    )
    add_text(
        s,
        MARGIN,
        Inches(4.0),
        Inches(7),
        Inches(0.6),
        "The simulation, eval & observability layer\nfor production voice AI.",
        size=16,
        color=MUTED,
    )
    hairline(s, MARGIN, Inches(5.0), Inches(2.5))
    add_text(s, MARGIN, Inches(5.2), Inches(7), Inches(0.3), "Bharat Bhavnasi  ·  Founder & CEO", size=13, color=INK)
    add_text(s, MARGIN, Inches(5.55), Inches(7), Inches(0.3), "Seed Round  ·  Bharat@Simkit.ai", size=12, color=MUTED)

    # soft callout chip
    pill(s, MARGIN, Inches(6.2), Inches(4.4), Inches(0.45), CALL_YELLOW)
    add_text(
        s,
        MARGIN + Inches(0.2),
        Inches(6.25),
        Inches(4.0),
        Inches(0.35),
        "✦  Test every caller before they call you",
        size=12,
        bold=True,
        color=INK,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # ══════════════════════════════════════════════════════════
    # 02 FOUNDER — doodle avatar + fit
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "01  ·  FOUNDER")
    headline(s, "Built for voice agents\nby someone who breaks them.", size=32, width=Inches(7.5))

    # big doodle avatar card
    card(s, Inches(9.0), Inches(1.0), Inches(3.5), Inches(5.3), fill=SURFACE)
    add_pic(s, doodle("03"), Inches(9.35), Inches(1.4), Inches(2.8), Inches(2.8))
    add_text(s, Inches(9.2), Inches(4.4), Inches(3.1), Inches(0.35), "Bharat Bhavnasi", size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.2), Inches(4.8), Inches(3.1), Inches(0.3), "Founder & CEO · SF", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    # mini persona row under avatar
    for i, n in enumerate(["01", "40", "12"]):
        add_pic(s, doodle(n), Inches(9.45) + i * Inches(0.95), Inches(5.35), Inches(0.8), Inches(0.8))

    add_text(s, MARGIN, Inches(2.9), Inches(7.5), Inches(0.3), "Why me", size=12, bold=True, color=MUTED)
    bullets = [
        "•  15 years shipping enterprise products",
        "•  4 years deep in AI & voice-agent orchestration",
        "•  3 patents filed: n-Brain, cross-layer context, voice-driven web",
        "•  Advisor to AI startups; helped founders ship v1 & raise seed",
    ]
    add_lines(s, MARGIN, Inches(3.3), Inches(7.5), Inches(2.0), bullets, size=15, color=INK)

    pill(s, MARGIN, Inches(5.7), Inches(7.6), Inches(0.7), CALL_ORANGE)
    add_text(
        s,
        MARGIN + Inches(0.25),
        Inches(5.8),
        Inches(7.1),
        Inches(0.5),
        "Founder–market fit: I build the tooling I needed while building agents.",
        size=13,
        bold=True,
        color=INK,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    footer(s, "02")

    # ══════════════════════════════════════════════════════════
    # 03 PROBLEM — icons + doodles + stats
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "02  ·  PROBLEM")
    add_pic(s, icon("problem"), Inches(11.7), Inches(0.35), Inches(0.9), Inches(0.9))
    headline(s, "Building a voice agent is easy.\nMaking it production-ready is brutal.", size=30)
    sub(
        s,
        "Agents break on accents, interruptions, latency, missing context, and compliance — and teams find out when customers do.",
        top=Inches(2.2),
    )

    # persona chaos row
    for i, n in enumerate(["05", "22", "27", "35", "48", "08"]):
        add_pic(s, doodle(n), MARGIN + i * Inches(1.15), Inches(3.0), Inches(1.0), Inches(1.0))
    add_text(s, MARGIN + Inches(7.1), Inches(3.25), Inches(5), Inches(0.5), "Real callers. Real edge cases. Real breakage.", size=13, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)

    stats = [
        ("40%", "agentic AI projects\ncanceled by 2027", CALL_PINK),
        ("15%", "teams reach strong\neval coverage", CALL_YELLOW),
        ("~40%", "voice failures missed\nby text-only tools", CALL_ORANGE),
    ]
    for i, (num, label, tint) in enumerate(stats):
        left = MARGIN + i * Inches(4.0)
        card(s, left, Inches(4.3), Inches(3.75), Inches(2.15), fill=tint)
        add_text(s, left + Inches(0.3), Inches(4.5), Inches(3.15), Inches(0.55), num, size=32, bold=True, color=INK)
        add_text(s, left + Inches(0.3), Inches(5.2), Inches(3.15), Inches(0.9), label, size=13, color=MUTED)
    footer(s, "03")

    # ══════════════════════════════════════════════════════════
    # 04 WHY NOW
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "03  ·  WHY NOW")
    add_pic(s, icon("phone"), Inches(11.7), Inches(0.35), Inches(0.9), Inches(0.9))
    headline(s, "Voice is becoming the default\ninterface to AI.", size=32)

    metrics = [
        ("$7B+", "into voice AI\nin Q1 2026 alone", CALL_YELLOW),
        ("40%", "of service calls\nAI-handled by 2027", CALL_GREEN),
        ("87.5%", "of builders shipping\nvoice agents now", CALL_ORANGE),
        ("$100B+", "voice + conversational AI\nby early 2030s", CALL_PINK),
    ]
    for i, (num, label, tint) in enumerate(metrics):
        left = MARGIN + i * Inches(3.1)
        card(s, left, Inches(2.9), Inches(2.9), Inches(2.55), fill=tint)
        add_text(s, left + Inches(0.25), Inches(3.15), Inches(2.4), Inches(0.55), num, size=24, bold=True, color=INK)
        add_text(s, left + Inches(0.25), Inches(3.85), Inches(2.4), Inches(1.2), label, size=13, color=MUTED)

    pill(s, MARGIN, Inches(5.8), Inches(11.9), Inches(0.7), SURFACE)
    add_text(
        s,
        MARGIN + Inches(0.3),
        Inches(5.9),
        Inches(11.3),
        Inches(0.5),
        "In ~2 years, every business phone number will be a voice agent — each one needs to be tested before it answers.",
        size=14,
        color=INK,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    footer(s, "04")

    # ══════════════════════════════════════════════════════════
    # 05 SOLUTION — Simulate / Evaluate / Observe + doodles
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "04  ·  SOLUTION")
    headline(s, "One harness to test\nevery voice agent.", size=32)
    sub(s, "Not a model wrapper — the simulation layer every voice-AI team needs to ship with confidence.")

    pillars = [
        ("simulate", "Simulate", "Thousands of synthetic callers — accents, noise, interruptions, emotion.", CALL_YELLOW, ["01", "40", "05"]),
        ("evaluate", "Evaluate", "LLM-judge + deterministic scorers for quality, task success, latency, safety.", CALL_ORANGE, ["12", "18"]),
        ("observe", "Observe", "Live monitoring — production failures become new test cases automatically.", CALL_GREEN, ["30", "44"]),
    ]
    for i, (ic, title, body, tint, dns) in enumerate(pillars):
        left = MARGIN + i * Inches(4.05)
        card(s, left, Inches(3.05), Inches(3.85), Inches(3.4), fill=WHITE)
        # tint header strip
        pill(s, left + Inches(0.15), Inches(3.2), Inches(3.55), Inches(0.95), tint)
        add_pic(s, icon(ic), left + Inches(0.3), Inches(3.32), Inches(0.7), Inches(0.7))
        add_text(s, left + Inches(1.15), Inches(3.4), Inches(2.3), Inches(0.5), title, size=18, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + Inches(0.3), Inches(4.35), Inches(3.25), Inches(1.3), body, size=13, color=MUTED)
        for j, n in enumerate(dns[:3]):
            add_pic(s, doodle(n), left + Inches(0.3) + j * Inches(0.7), Inches(5.75), Inches(0.55), Inches(0.55))
    footer(s, "05")

    # ══════════════════════════════════════════════════════════
    # 06 PRODUCT — icon grid + persona strip
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "05  ·  PRODUCT")
    add_pic(s, icon("api"), Inches(11.7), Inches(0.35), Inches(0.9), Inches(0.9))
    headline(s, "Integrates with any voice agent.", size=32)

    features = [
        ("simulate", "Voice-to-voice sim", "Agent-vs-agent calls, end to end, at scale", CALL_YELLOW),
        ("phone", "Live phone testing", "Dial a deployed agent; validate production-ready", CALL_ORANGE),
        ("evaluate", "Chat agent testing", "LLM-to-LLM evaluation for text agents too", CALL_GREEN),
        ("api", "API drop-in", "LiveKit, Vapi, Retell, ElevenLabs, Pipecat", CALL_GRAY),
        ("idea", "Test personas", "Author personas that probe every edge case", CALL_PINK),
        ("target", "CI/CD loop", "Personas → Simulate → Score → Observe", CALL_YELLOW),
    ]
    for i, (ic, title, body, tint) in enumerate(features):
        row, col = divmod(i, 3)
        left = MARGIN + col * Inches(4.05)
        top = Inches(2.25) + row * Inches(2.15)
        card(s, left, top, Inches(3.85), Inches(1.95), fill=WHITE)
        pill(s, left + Inches(0.2), top + Inches(0.25), Inches(0.7), Inches(0.7), tint)
        add_pic(s, icon(ic), left + Inches(0.28), top + Inches(0.33), Inches(0.55), Inches(0.55))
        add_text(s, left + Inches(1.1), top + Inches(0.3), Inches(2.5), Inches(0.4), title, size=15, bold=True, color=INK)
        add_text(s, left + Inches(1.1), top + Inches(0.8), Inches(2.5), Inches(0.85), body, size=12, color=MUTED)
    footer(s, "06")

    # ══════════════════════════════════════════════════════════
    # 07 MARKET
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "06  ·  MARKET")
    add_pic(s, icon("market"), Inches(11.7), Inches(0.35), Inches(0.9), Inches(0.9))
    headline(s, "A reliability layer on\na fast-growing wave.", size=32)
    sub(s, "Voice + conversational AI ~$22B today → $100B+ by early 2030s. Our slice: testing & reliability.")

    markets = [
        ("TAM", "$12B", "AI agent quality & reliability by 2030", CALL_YELLOW),
        ("SAM", "$3.5B", "Voice-first testing · BFSI, health, telecom", CALL_ORANGE),
        ("SOM", "$120M", "Year-5 ARR · ~2,400 teams @ ~$50K ACV", CALL_GREEN),
    ]
    for i, (tag, num, label, tint) in enumerate(markets):
        left = MARGIN + i * Inches(4.05)
        card(s, left, Inches(3.2), Inches(3.85), Inches(3.0), fill=WHITE)
        pill(s, left + Inches(0.3), Inches(3.45), Inches(1.0), Inches(0.35), tint)
        add_text(s, left + Inches(0.4), Inches(3.48), Inches(0.8), Inches(0.3), tag, size=11, bold=True, color=INK)
        add_text(s, left + Inches(0.3), Inches(4.05), Inches(3.25), Inches(0.7), num, size=36, bold=True, color=INK)
        add_text(s, left + Inches(0.3), Inches(4.9), Inches(3.25), Inches(0.9), label, size=14, color=MUTED)
    footer(s, "07")

    # ══════════════════════════════════════════════════════════
    # 08 WHY SIMKIT WINS
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "07  ·  WHY SIMKIT WINS")
    add_pic(s, icon("shield"), Inches(11.7), Inches(0.35), Inches(0.9), Inches(0.9))
    headline(s, "Everyone can score.\nNo one is voice-native and unified.", size=30)

    wins = [
        ("phone", "Voice-first", "Not a text retrofit — accents, barge-in, TTFA, TTS.", CALL_YELLOW),
        ("target", "One loop", "Sim + eval + observability — not three tools.", CALL_ORANGE),
        ("api", "Partner, not rival", "We don’t compete with agent builders — we arm them.", CALL_GREEN),
        ("shield", "Deep moat", "3 patents filed in voice-agent orchestration.", CALL_PINK),
    ]
    for i, (ic, title, body, tint) in enumerate(wins):
        left = MARGIN + (i % 2) * Inches(6.05)
        top = Inches(2.85) + (i // 2) * Inches(1.8)
        card(s, left, top, Inches(5.8), Inches(1.55), fill=WHITE)
        pill(s, left + Inches(0.25), top + Inches(0.35), Inches(0.8), Inches(0.8), tint)
        add_pic(s, icon(ic), left + Inches(0.35), top + Inches(0.45), Inches(0.6), Inches(0.6))
        add_text(s, left + Inches(1.3), top + Inches(0.3), Inches(4.2), Inches(0.4), title, size=16, bold=True, color=INK)
        add_text(s, left + Inches(1.3), top + Inches(0.8), Inches(4.2), Inches(0.5), body, size=13, color=MUTED)
    footer(s, "08")

    # ══════════════════════════════════════════════════════════
    # 09 TRACTION & MODEL
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "08  ·  TRACTION & MODEL")
    add_pic(s, icon("rocket"), Inches(11.7), Inches(0.35), Inches(0.9), Inches(0.9))
    headline(s, "Builders already want this.", size=32)

    proof = [
        ("Working product", "Product-0 built & demoable; production-ready path ~5 months.", CALL_YELLOW),
        ("Top validation", "Demoed to CTOs of LiveKit & ElevenLabs.", CALL_ORANGE),
        ("100+ practitioners", "Strong pull across nearly every voice provider network.", CALL_GREEN),
    ]
    for i, (title, body, tint) in enumerate(proof):
        left = MARGIN + i * Inches(4.05)
        card(s, left, Inches(2.2), Inches(3.85), Inches(1.75), fill=tint)
        add_text(s, left + Inches(0.3), Inches(2.4), Inches(3.25), Inches(0.35), title, size=14, bold=True, color=INK)
        add_text(s, left + Inches(0.3), Inches(2.9), Inches(3.25), Inches(0.8), body, size=13, color=MUTED)

    add_text(s, MARGIN, Inches(4.2), Inches(6), Inches(0.3), "Open-core, usage-based", size=13, bold=True, color=MUTED)

    tiers = [
        ("Open source", "Free", "Simulator wedge"),
        ("Team", "$1.5–5K/mo", "~$25–60K ACV"),
        ("Enterprise", "$100K+/yr", "SSO, VPC, SLAs"),
    ]
    for i, (name, price, note) in enumerate(tiers):
        left = MARGIN + i * Inches(4.05)
        card(s, left, Inches(4.65), Inches(3.85), Inches(1.7), fill=WHITE)
        add_text(s, left + Inches(0.3), Inches(4.85), Inches(3.25), Inches(0.3), name, size=14, bold=True, color=INK)
        add_text(s, left + Inches(0.3), Inches(5.25), Inches(3.25), Inches(0.35), price, size=20, bold=True, color=INK)
        add_text(s, left + Inches(0.3), Inches(5.7), Inches(3.25), Inches(0.3), note, size=12, color=MUTED)
    footer(s, "09")

    # ══════════════════════════════════════════════════════════
    # 10 THE ASK
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    fill_slide(s, WHITE)
    kicker(s, "09  ·  THE ASK")
    add_pic(s, icon("rocket"), Inches(11.7), Inches(0.35), Inches(0.9), Inches(0.9))
    headline(s, "Raising $4M for 10%", size=40)
    sub(s, "to make every voice agent production-ready — ~20 months of runway.", top=Inches(1.7))

    # doodle row for energy
    for i, n in enumerate(["01", "40", "12", "18", "30", "44"]):
        add_pic(s, doodle(n), MARGIN + i * Inches(0.85), Inches(2.4), Inches(0.7), Inches(0.7))

    uses = [
        ("Product", "Ship v1 in ~5 months — simulator, evals & live observability", CALL_YELLOW),
        ("Team", "5–8 engineers + 3 GTM; hire a marketing co-founder", CALL_ORANGE),
        ("Customers", "Land 2–5 enterprise design partners in year one", CALL_GREEN),
    ]
    for i, (title, body, tint) in enumerate(uses):
        top = Inches(3.35) + i * Inches(0.85)
        card(s, MARGIN, top, Inches(11.9), Inches(0.72), fill=tint)
        add_text(s, MARGIN + Inches(0.35), top + Inches(0.18), Inches(2.0), Inches(0.4), title, size=15, bold=True, color=INK)
        add_text(s, MARGIN + Inches(2.5), top + Inches(0.18), Inches(9.0), Inches(0.4), body, size=14, color=INK)

    add_text(
        s,
        MARGIN,
        Inches(6.1),
        Inches(11.5),
        Inches(0.3),
        "Interested? Live demo anytime.  ·  Bharat Bhavnasi  ·  Bharat@simkit.ai  ·  github.com/bvsbharat",
        size=12,
        color=MUTED,
    )
    add_text(
        s,
        MARGIN,
        Inches(6.45),
        Inches(11.5),
        Inches(0.35),
        "Every voice agent deserves a soundcheck.",
        size=16,
        bold=True,
        color=INK,
    )
    footer(s, "10")

    # ══════════════════════════════════════════════════════════
    # APPENDIX DIVIDER
    # ══════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    fill_slide(s, SURFACE)
    for i, n in enumerate(["01", "05", "12", "18", "30", "40"]):
        add_pic(s, doodle(n), Inches(3.5) + i * Inches(1.1), Inches(2.0), Inches(0.95), Inches(0.95))
    add_text(s, MARGIN, Inches(3.4), Inches(11.5), Inches(0.4), "APPENDIX", size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(3.9), Inches(11.5), Inches(0.7), "Detail for diligence.", size=36, bold=True, color=INK, align=PP_ALIGN.CENTER)

    # ── A1 HOW IT WORKS ───────────────────────────────────────
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "A1  ·  HOW IT WORKS")
    headline(s, "A continuous quality loop, wired into CI/CD.", size=28)
    steps = [
        ("idea", "1 · Personas", "Define caller & red-team personas as code", CALL_YELLOW, "01"),
        ("simulate", "2 · Simulate", "Thousands of real audio calls per release", CALL_ORANGE, "40"),
        ("evaluate", "3 · Score", "Grade every turn; localize the failing step", CALL_GREEN, "12"),
        ("observe", "4 · Observe", "Prod failures become new test cases", CALL_PINK, "30"),
    ]
    for i, (ic, title, body, tint, dn) in enumerate(steps):
        left = MARGIN + i * Inches(3.1)
        card(s, left, Inches(2.6), Inches(2.95), Inches(3.6), fill=WHITE)
        pill(s, left + Inches(0.25), Inches(2.85), Inches(0.7), Inches(0.7), tint)
        add_pic(s, icon(ic), left + Inches(0.33), Inches(2.93), Inches(0.55), Inches(0.55))
        add_text(s, left + Inches(0.25), Inches(3.75), Inches(2.45), Inches(0.5), title, size=15, bold=True, color=INK)
        add_text(s, left + Inches(0.25), Inches(4.35), Inches(2.45), Inches(1.0), body, size=13, color=MUTED)
        add_pic(s, doodle(dn), left + Inches(0.95), Inches(5.5), Inches(0.9), Inches(0.9))
    footer(s, "A1")

    # ── A2 METRICS ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "A2  ·  METRICS")
    headline(s, "Signals that make a voice agent production-ready.", size=28)
    metrics = [
        ("TTFA", "Felt latency that makes or breaks UX", CALL_YELLOW),
        ("TTFB / TTFT", "STT → LLM → TTS path latency", CALL_ORANGE),
        ("Task completion", "Did the agent hit the caller’s goal?", CALL_GREEN),
        ("Tool-call accuracy", "Correct calls, args, recovery", CALL_PINK),
        ("Interruptions", "Barge-in, silence, overlap, ASR", CALL_GRAY),
        ("Safety", "PII, compliance, instruction fidelity", CALL_YELLOW),
    ]
    for i, (title, body, tint) in enumerate(metrics):
        row, col = divmod(i, 3)
        left = MARGIN + col * Inches(4.05)
        top = Inches(2.4) + row * Inches(2.05)
        card(s, left, top, Inches(3.85), Inches(1.85), fill=tint)
        add_text(s, left + Inches(0.3), top + Inches(0.35), Inches(3.25), Inches(0.4), title, size=15, bold=True, color=INK)
        add_text(s, left + Inches(0.3), top + Inches(0.9), Inches(3.25), Inches(0.7), body, size=13, color=MUTED)
    footer(s, "A2")

    # ── A3 FULL-STACK ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "A3  ·  FULL-STACK STRESS TEST")
    headline(s, "v1.0 attacks every layer of the voice stack.", size=28)
    layers = [
        ("STT · Speech-in", "Accents, noise, cross-talk, disfluencies, code-switching", CALL_YELLOW, "05"),
        ("LLM · Reasoning", "Context, tool-calls, hallucination, PII & compliance", CALL_ORANGE, "12"),
        ("TTS · Speech-out", "Barge-in, TTFA/TTFB, naturalness & emotion", CALL_GREEN, "40"),
    ]
    for i, (title, body, tint, dn) in enumerate(layers):
        top = Inches(2.55) + i * Inches(1.25)
        card(s, MARGIN, top, Inches(11.9), Inches(1.1), fill=tint)
        add_pic(s, doodle(dn), MARGIN + Inches(0.25), top + Inches(0.2), Inches(0.7), Inches(0.7))
        add_text(s, MARGIN + Inches(1.2), top + Inches(0.25), Inches(3.5), Inches(0.6), title, size=16, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, MARGIN + Inches(5.0), top + Inches(0.25), Inches(7.2), Inches(0.6), body, size=14, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "A3")

    # ── A4 VALUATION ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "A4  ·  VALUATION LANDSCAPE")
    headline(s, "Capital is pouring into this category.", size=28)
    cos = [
        ("ElevenLabs", "$11B · $500M ARR", CALL_YELLOW),
        ("Deepgram", "$1.3B · Series C", CALL_ORANGE),
        ("PolyAI", "$750M · Series D", CALL_GREEN),
        ("Vapi", "~$500M · 1B+ calls", CALL_PINK),
    ]
    for i, (name, note, tint) in enumerate(cos):
        left = MARGIN + i * Inches(3.1)
        card(s, left, Inches(2.7), Inches(2.95), Inches(2.5), fill=tint)
        add_text(s, left + Inches(0.25), Inches(3.1), Inches(2.45), Inches(0.5), name, size=16, bold=True, color=INK)
        add_text(s, left + Inches(0.25), Inches(3.75), Inches(2.45), Inches(0.9), note, size=14, color=MUTED)
    add_text(
        s,
        MARGIN,
        Inches(5.6),
        Inches(11.5),
        Inches(0.6),
        "Voice-agent testing & evaluation — Simkit’s category — is being funded now.",
        size=15,
        color=INK,
    )
    footer(s, "A4")

    # ── A5 DEFENSIBILITY ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "A5  ·  DEFENSIBILITY")
    add_pic(s, icon("shield"), Inches(11.7), Inches(0.35), Inches(0.9), Inches(0.9))
    headline(s, "A moat built on deep voice-agent research.", size=28)
    moats = [
        ("Dual / n-Brain", "Parallel inner + conversational + proactive brain for low-latency complex flows", CALL_YELLOW),
        ("Cross-layer context", "STT · LLM · TTS share emotion, tone, intent — ~800ms responses", CALL_ORANGE),
        ("Voice-driven web", "A voice agent that drives real web operations end-to-end", CALL_GREEN),
        ("Faster inference", "Flash + reasoning model coordination for snappier agents", CALL_PINK),
    ]
    for i, (title, body, tint) in enumerate(moats):
        row, col = divmod(i, 2)
        left = MARGIN + col * Inches(6.05)
        top = Inches(2.5) + row * Inches(1.9)
        card(s, left, top, Inches(5.8), Inches(1.7), fill=tint)
        add_text(s, left + Inches(0.35), top + Inches(0.3), Inches(5.1), Inches(0.4), title, size=16, bold=True, color=INK)
        add_text(s, left + Inches(0.35), top + Inches(0.85), Inches(5.1), Inches(0.6), body, size=13, color=MUTED)
    footer(s, "A5")

    # ── A6 ABOUT FOUNDER ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "A6  ·  ABOUT THE FOUNDER")
    headline(s, "A builder who turns hard agent problems into products.", size=28)
    stats = [
        ("15 yrs", "enterprise products", CALL_YELLOW),
        ("4 yrs", "AI & voice agents", CALL_ORANGE),
        ("50+", "apps built", CALL_GREEN),
        ("3", "patents filed", CALL_PINK),
    ]
    for i, (num, label, tint) in enumerate(stats):
        left = MARGIN + i * Inches(3.1)
        card(s, left, Inches(2.6), Inches(2.95), Inches(2.4), fill=tint)
        add_text(s, left + Inches(0.25), Inches(2.95), Inches(2.45), Inches(0.6), num, size=28, bold=True, color=INK)
        add_text(s, left + Inches(0.25), Inches(3.7), Inches(2.45), Inches(0.7), label, size=14, color=MUTED)
    add_text(
        s,
        MARGIN,
        Inches(5.5),
        Inches(11.5),
        Inches(0.7),
        "I'm not competing with agent builders — I'm giving every one of them the tooling to ship quality.",
        size=16,
        bold=True,
        color=INK,
    )
    footer(s, "A6")

    # ── A7 MILESTONES ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    fill_slide(s)
    kicker(s, "A7  ·  12–18 MONTH MILESTONES")
    add_pic(s, icon("rocket"), Inches(11.7), Inches(0.35), Inches(0.9), Inches(0.9))
    headline(s, "What $4M unlocks.", size=32)
    milestones = [
        ("01", "Open-source simulator becomes the default way teams test voice agents", CALL_YELLOW),
        ("02", "2–5 enterprise design partners → first reference customers & ARR", CALL_ORANGE),
        ("03", "Partnerships with leading voice providers (LiveKit, ElevenLabs, and peers)", CALL_GREEN),
        ("04", "v1: simulator + evals + live observability in production", CALL_PINK),
    ]
    for i, (num, body, tint) in enumerate(milestones):
        top = Inches(2.45) + i * Inches(1.0)
        card(s, MARGIN, top, Inches(11.9), Inches(0.85), fill=tint)
        add_text(s, MARGIN + Inches(0.35), top + Inches(0.22), Inches(0.7), Inches(0.4), num, size=18, bold=True, color=INK)
        add_text(s, MARGIN + Inches(1.2), top + Inches(0.22), Inches(10.2), Inches(0.45), body, size=15, color=INK)

    footer(s, "A7")

    prs.save(OUT_LOCAL)
    print(f"Saved {OUT_LOCAL}")
    print(f"Slides: {len(prs.slides)}")
    for dest in (OUT_DESKTOP, OUT_DOWNLOADS):
        try:
            prs.save(dest)
            print(f"Saved {dest}")
        except Exception as e:
            print(f"skip {dest}: {e}")


if __name__ == "__main__":
    build()
